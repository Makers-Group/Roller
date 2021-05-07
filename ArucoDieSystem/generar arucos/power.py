from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_CONNECTOR
imagelist = []
def generar_power(inicio, fin, nombre):
    arucosx = 0
    arucosy = 0
#    inicio = 0
    prs = Presentation("template.pptx")
    blank_slide_layout = prs.slide_layouts[6]
#    slide = prs.slides.add_slide(blank_slide_layout)
    slide = prs.slides[0]
    x = 0.6
    y = 0.6
    for i in range(inicio,fin):
            img = 'id_' + str(i) + '.png'
            imagelist.append(img)
    for img in imagelist:
        txBox = slide.shapes.add_textbox(Cm(x), Cm(y-.7), Cm(5), Cm(3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Id: " + str(inicio)
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(16)
        font.bold = True
        pic = slide.shapes.add_picture(img, Cm(x), Cm(y) ,height=Cm(4.2))
        x += 5.4
        arucosx += 1
        inicio += 1
        line1=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x-.6), Cm(0), Cm(x-.6), Cm(65))
        if arucosx == 16:
            x = 0.6
            arucosx = 0
            y += 5.4
            arucosy += 1
            line2=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(0), Cm(y-.6), Cm(87), Cm(y-.6))
        if arucosy == 12:
            arucosy = 0
            slide = prs.slides.add_slide(blank_slide_layout)
            x = 0.6
            y = 0.6


    prs.save(nombre + '.pptx')

generar_power(300,701,"Arucos")
    

    
